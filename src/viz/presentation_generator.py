#!/usr/bin/env python
"""
Montana State Agency Presentation Generator

Generates 30-minute presentations for each Montana state agency based on their
website content and supporting materials. Creates both Markdown and PowerPoint formats.

Usage:
    python visualizations/presentation_generator.py                    # Generate all agency presentations
    python visualizations/presentation_generator.py --agency commerce  # Generate specific agency
    python visualizations/presentation_generator.py --md-only          # Generate markdown only
"""

import argparse
import json
import re
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from collections import Counter
from datetime import datetime
import sys

# Add parent directory to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent.parent))

from src.network.graph_builder import GraphBuilder
from src.network.persistence import GraphPersistence
from src.network.rag_retriever import GraphRAGRetriever


class AgencyContentAnalyzer:
    """Analyzes agency content to extract key presentation points"""
    
    def __init__(self, knowledge_dir: str = "../knowledge"):
        self.knowledge_dir = Path(knowledge_dir)
        self.graph = None
        self.retriever = None
        
    def load_network(self, graph_path: Optional[str] = None):
        """Load existing knowledge network or build new one"""
        if graph_path and Path(graph_path).exists():
            print(f"Loading existing graph from {graph_path}...")
            # Extract just the filename from the path
            graph_file = Path(graph_path)
            output_dir = str(graph_file.parent)
            filename = graph_file.name
            persistence = GraphPersistence(output_dir=output_dir)
            self.graph = persistence.load_pickle(filename)
        else:
            print("Building knowledge graph...")
            builder = GraphBuilder(knowledge_dir=str(self.knowledge_dir))
            self.graph = builder.build_graph(verbose=True)
        
        if self.graph:
            self.retriever = GraphRAGRetriever(self.graph)
        
    def get_agency_documents(self, agency_folder: str) -> List[Dict]:
        """Get all documents for a specific agency"""
        agency_docs = []
        
        for node_id, node in self.graph.nodes.items():
            if node.agency == agency_folder:
                agency_docs.append({
                    'id': node_id,
                    'title': node.title,
                    'type': node.node_type.value if hasattr(node.node_type, 'value') else str(node.node_type),
                    'keywords': node.keywords,
                    'topics': node.topics,
                    'url': node.source_url,
                    'pagerank': node.pagerank_score,
                    'content_preview': self._get_content_preview(node)
                })
        
        # Sort by importance (pagerank)
        agency_docs.sort(key=lambda x: x['pagerank'], reverse=True)
        return agency_docs
    
    def _get_content_preview(self, node) -> str:
        """Get content preview from chunks"""
        chunks = node.chunks if hasattr(node, 'chunks') and node.chunks else []
        if chunks:
            first_chunk = chunks[0]
            content = first_chunk.content if hasattr(first_chunk, 'content') else str(first_chunk)
            return content[:500]
        return ""
    
    def extract_services(self, agency_folder: str) -> Tuple[List[Dict], List[Tuple]]:
        """Extract services provided by the agency"""
        # Search for service-related content
        results = self.retriever.search_by_agency(
            agency=agency_folder,
            query="services programs offerings citizens residents benefits applications",
            top_k=50
        )
        
        services = []
        service_keywords = []
        
        for doc in results:
            node_id = doc.get('node_id')
            if not node_id:
                continue
            
            # Get chunks for this document
            node = self.graph.nodes.get(node_id)
            if not node:
                continue
            
            chunks = node.chunks if hasattr(node, 'chunks') else []
            keywords = node.keywords if hasattr(node, 'keywords') else []
            
            for chunk in chunks[:3]:  # Limit chunks per document
                content = chunk.content if hasattr(chunk, 'content') else str(chunk)
                # Look for service indicators
                if any(word in content.lower() for word in ['service', 'program', 'benefit', 'application', 'assistance']):
                    services.append({
                        'title': doc.get('title', ''),
                        'content': content,
                        'url': doc.get('source_url', '')
                    })
                    service_keywords.extend(keywords)
        
        # Get most common service areas
        common_keywords = Counter(service_keywords).most_common(10)
        
        return services, common_keywords
    
    def extract_mission_overview(self, agency_folder: str, agency_name: str) -> Dict:
        """Extract agency mission and overview"""
        # Search for mission/about content
        results = self.retriever.search_by_agency(
            agency=agency_folder,
            query="mission vision about overview purpose role responsibility",
            top_k=10
        )
        
        mission_content = []
        key_docs = []
        
        for doc in results:
            node_id = doc.get('node_id')
            title = doc.get('title', '')
            
            if any(word in title.lower() for word in ['about', 'mission', 'overview', 'index']):
                # Get chunks for this document
                if node_id and node_id in self.graph.nodes:
                    node = self.graph.nodes[node_id]
                    chunks = node.chunks if hasattr(node, 'chunks') else []
                    for chunk in chunks[:2]:
                        content = chunk.content if hasattr(chunk, 'content') else str(chunk)
                        mission_content.append(content)
                
                key_docs.append(doc)
        
        return {
            'agency_name': agency_name,
            'folder': agency_folder,
            'mission_text': ' '.join(mission_content[:3]) if mission_content else f"Montana {agency_name}",
            'key_documents': key_docs[:5]
        }
    
    def categorize_services(self, services: List[Dict], keywords: List[Tuple]) -> Dict[str, List]:
        """Categorize services into logical groups"""
        categories = {
            'Citizen Services': [],
            'Business Services': [],
            'Government Services': [],
            'Programs & Benefits': [],
            'Regulations & Compliance': [],
            'Resources & Information': []
        }
        
        # Simple categorization based on keywords and content
        for service in services:
            content_lower = service['content'].lower()
            title_lower = service['title'].lower()
            
            categorized = False
            
            # Citizen services
            if any(word in content_lower or word in title_lower for word in 
                   ['citizen', 'resident', 'public', 'individual', 'family', 'personal']):
                categories['Citizen Services'].append(service)
                categorized = True
            
            # Business services
            if any(word in content_lower or word in title_lower for word in 
                   ['business', 'commercial', 'industry', 'employer', 'license', 'permit']):
                categories['Business Services'].append(service)
                categorized = True
            
            # Programs & Benefits
            if any(word in content_lower or word in title_lower for word in 
                   ['program', 'benefit', 'assistance', 'grant', 'fund', 'support']):
                categories['Programs & Benefits'].append(service)
                categorized = True
            
            # Regulations
            if any(word in content_lower or word in title_lower for word in 
                   ['regulation', 'compliance', 'requirement', 'rule', 'law', 'policy']):
                categories['Regulations & Compliance'].append(service)
                categorized = True
            
            # Resources
            if any(word in content_lower or word in title_lower for word in 
                   ['resource', 'information', 'guide', 'form', 'document', 'report']):
                categories['Resources & Information'].append(service)
                categorized = True
            
            # Government services (default or explicit)
            if not categorized or any(word in content_lower or word in title_lower for word in 
                   ['government', 'agency', 'department', 'division', 'office']):
                categories['Government Services'].append(service)
        
        # Remove empty categories and limit items
        return {k: v[:8] for k, v in categories.items() if v}


class MarkdownPresentationGenerator:
    """Generates markdown presentations suitable for conversion to PowerPoint"""
    
    def __init__(self, output_dir: str = "../presentations"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
    
    def generate_presentation(self, agency_data: Dict, services_data: Dict) -> str:
        """Generate complete markdown presentation"""
        
        agency_name = agency_data['agency_name']
        folder = agency_data['folder']
        
        # Create agency-specific folder
        agency_dir = self.output_dir / folder
        agency_dir.mkdir(exist_ok=True)
        
        # Generate markdown content
        md_content = self._create_markdown_slides(agency_name, agency_data, services_data)
        
        # Save markdown
        md_file = agency_dir / f"{folder}_presentation.md"
        md_file.write_text(md_content, encoding='utf-8')
        
        print(f"✓ Generated markdown: {md_file}")
        return str(md_file)
    
    def _create_markdown_slides(self, agency_name: str, agency_data: Dict, services_data: Dict) -> str:
        """Create markdown presentation content"""
        
        slides = []
        
        # Slide 1: Title Slide
        slides.append(f"""---

# {agency_name}

## Services for Montana Citizens

State of Montana

{datetime.now().strftime('%B %Y')}

---
""")
        
        # Slide 2: Agenda
        slides.append("""---

## Presentation Agenda

1. Agency Overview & Mission
2. Key Services & Programs
3. Citizen Services
4. Business Services
5. Programs & Benefits
6. Resources & Support
7. How to Access Services
8. Contact Information

---
""")
        
        # Slide 3: Mission & Overview
        mission_text = agency_data.get('mission_text', '')[:600]
        slides.append(f"""---

## Agency Mission & Overview

{mission_text}

### Our Role

Serving Montana citizens through dedicated programs and services that support our communities and economy.

---
""")
        
        # Slides 4-N: Service Categories
        categories = services_data.get('categories', {})
        
        for category_name, services in categories.items():
            if not services:
                continue
                
            slides.append(f"""---

## {category_name}

""")
            
            # Add up to 6 services per category
            service_items = []
            seen_titles = set()
            
            for service in services[:6]:
                title = service['title']
                if title in seen_titles:
                    continue
                seen_titles.add(title)
                
                # Clean and summarize content
                content = service['content'][:200].strip()
                # Remove extra whitespace
                content = re.sub(r'\s+', ' ', content)
                
                service_items.append(f"- **{title}**\n  - {content}...")
            
            slides[-1] += '\n'.join(service_items)
            slides[-1] += "\n\n---\n"
        
        # Slide: Key Service Areas (from keywords)
        keywords = services_data.get('keywords', [])[:10]
        if keywords:
            slides.append("""---

## Key Service Areas

""")
            keyword_list = '\n'.join([f"- {word.title()} ({count} references)" for word, count in keywords])
            slides[-1] += keyword_list + "\n\n---\n"
        
        # Slide: How to Access Services
        slides.append(f"""---

## How to Access Our Services

### Online
- Visit our website for information and applications
- Many services available online 24/7

### In Person
- Visit our offices during business hours
- Schedule appointments for specialized services

### By Phone
- Contact our main office
- Speak with service representatives

### By Mail
- Submit applications and documents
- Request information packets

---
""")
        
        # Slide: Contact Information
        folder = agency_data['folder']
        url = agency_data.get('url', f'https://{folder}.mt.gov/')
        
        slides.append(f"""---

## Contact Information

### {agency_name}

**Website:** {url}

**General Information:**
- Visit our website for current contact details
- Office hours and locations
- Service-specific contact information

**Stay Connected:**
- Subscribe to updates
- Follow us on social media
- Sign up for notifications

---
""")
        
        # Final Slide: Thank You
        slides.append("""---

## Thank You

### Questions?

We're here to serve Montana citizens.

**Visit us online or contact our office for assistance.**

---
""")
        
        return '\n'.join(slides)


class PowerPointConverter:
    """Converts markdown presentations to PowerPoint with Montana government theme"""
    
    MONTANA_COLORS = {
        'primary_blue': '#003F87',      # Montana state blue
        'secondary_blue': '#4A90E2',    # Lighter blue
        'dark_gray': '#333333',         # Dark gray for text
        'medium_gray': '#666666',       # Medium gray
        'light_gray': '#E5E5E5',        # Light gray for backgrounds
        'white': '#FFFFFF',
        'accent_blue': '#0066CC'        # Accent blue
    }
    
    def __init__(self):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.dml.color import RGBColor
            self.pptx = True
            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
            self.PP_ALIGN = PP_ALIGN
            self.MSO_ANCHOR = MSO_ANCHOR
            self.RGBColor = RGBColor
        except ImportError:
            self.pptx = False
            print("Warning: python-pptx not installed. PowerPoint generation will be skipped.")
            print("Install with: pip install python-pptx")
    
    def convert_markdown_to_pptx(self, md_file: str, output_file: str):
        """Convert markdown presentation to PowerPoint"""
        if not self.pptx:
            print(f"Skipping PowerPoint conversion for {md_file}")
            return None
        
        # Read markdown content
        md_content = Path(md_file).read_text(encoding='utf-8')
        
        # Parse slides
        slides_content = self._parse_markdown_slides(md_content)
        
        # Create PowerPoint
        prs = self.Presentation()
        prs.slide_width = self.Inches(10)
        prs.slide_height = self.Inches(7.5)
        
        for slide_data in slides_content:
            self._create_slide(prs, slide_data)
        
        # Save PowerPoint
        prs.save(output_file)
        print(f"✓ Generated PowerPoint: {output_file}")
        return output_file
    
    def _parse_markdown_slides(self, md_content: str) -> List[Dict]:
        """Parse markdown into slide data structures"""
        slides = []
        
        # Split by slide delimiter
        slide_texts = md_content.split('---\n')
        
        for slide_text in slide_texts:
            slide_text = slide_text.strip()
            if not slide_text:
                continue
            
            lines = slide_text.split('\n')
            
            # Detect slide type and parse content
            slide_data = {
                'type': 'content',
                'title': '',
                'subtitle': '',
                'content': [],
                'bullets': []
            }
            
            in_content = False
            current_section = []
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # Title (# heading)
                if line.startswith('# ') and not slide_data['title']:
                    slide_data['title'] = line[2:].strip()
                    slide_data['type'] = 'title'
                
                # Subtitle (## heading)
                elif line.startswith('## '):
                    if not slide_data['title']:
                        slide_data['title'] = line[3:].strip()
                    else:
                        slide_data['subtitle'] = line[3:].strip()
                
                # Bullet points
                elif line.startswith('- ') or line.startswith('* '):
                    slide_data['bullets'].append(line[2:].strip())
                
                # Other content
                else:
                    if not line.startswith('#'):
                        slide_data['content'].append(line)
            
            slides.append(slide_data)
        
        return slides
    
    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def _create_slide(self, prs, slide_data: Dict):
        """Create a PowerPoint slide from slide data"""
        
        if slide_data['type'] == 'title':
            # Title slide layout
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = slide_data['title']
            
            # Set title formatting
            title.text_frame.paragraphs[0].font.size = self.Pt(44)
            title.text_frame.paragraphs[0].font.bold = True
            rgb = self._hex_to_rgb(self.MONTANA_COLORS['primary_blue'])
            title.text_frame.paragraphs[0].font.color.rgb = self.RGBColor(*rgb)
            
            # Add subtitle content
            if slide_data['subtitle']:
                subtitle.text = slide_data['subtitle']
            elif slide_data['content']:
                subtitle.text = '\n'.join(slide_data['content'][:3])
            
            # Subtitle formatting
            for paragraph in subtitle.text_frame.paragraphs:
                paragraph.font.size = self.Pt(24)
                rgb = self._hex_to_rgb(self.MONTANA_COLORS['dark_gray'])
                paragraph.font.color.rgb = self.RGBColor(*rgb)
        
        else:
            # Content slide layout
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            
            # Set title
            title.text = slide_data['title']
            title.text_frame.paragraphs[0].font.size = self.Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
            rgb = self._hex_to_rgb(self.MONTANA_COLORS['primary_blue'])
            title.text_frame.paragraphs[0].font.color.rgb = self.RGBColor(*rgb)
            
            # Add content
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                # Add subtitle if present
                if slide_data['subtitle']:
                    p = text_frame.paragraphs[0]
                    p.text = slide_data['subtitle']
                    p.font.size = self.Pt(24)
                    p.font.bold = True
                    rgb = self._hex_to_rgb(self.MONTANA_COLORS['secondary_blue'])
                    p.font.color.rgb = self.RGBColor(*rgb)
                    p.space_after = self.Pt(12)
                
                # Add bullets
                for bullet in slide_data['bullets']:
                    p = text_frame.add_paragraph()
                    
                    # Handle nested bullets (marked with **)
                    if '**' in bullet:
                        # Bold text between **
                        parts = bullet.split('**')
                        for i, part in enumerate(parts):
                            if i % 2 == 1:  # Bold part
                                run = p.add_run()
                                run.text = part
                                run.font.bold = True
                            else:  # Normal part
                                if part:
                                    run = p.add_run()
                                    run.text = part
                    else:
                        p.text = bullet
                    
                    p.font.size = self.Pt(18)
                    p.level = 0
                    rgb = self._hex_to_rgb(self.MONTANA_COLORS['dark_gray'])
                    p.font.color.rgb = self.RGBColor(*rgb)
                    p.space_after = self.Pt(6)
        
        # Set slide background
        background = slide.background
        fill = background.fill
        fill.solid()
        rgb = self._hex_to_rgb(self.MONTANA_COLORS['white'])
        fill.fore_color.rgb = self.RGBColor(*rgb)


def main():
    parser = argparse.ArgumentParser(
        description="Generate presentations for Montana state agencies",
        formatter_class=argparse.RawDescriptionHelpFormatter
    )
    
    parser.add_argument(
        "--agency",
        help="Generate presentation for specific agency (folder name)"
    )
    
    parser.add_argument(
        "--md-only",
        action="store_true",
        help="Generate markdown only (skip PowerPoint conversion)"
    )
    
    parser.add_argument(
        "--graph-path",
        default="../src/network/exports/montana_knowledge.pkl",
        help="Path to existing knowledge graph (default: ../src/network/exports/montana_knowledge.pkl)"
    )
    
    parser.add_argument(
        "--agencies-file",
        default="../agencies.md",
        help="Path to agencies.md file (default: ../agencies.md)"
    )
    
    args = parser.parse_args()
    
    # Read agencies list
    agencies_path = Path(args.agencies_file)
    if not agencies_path.exists():
        print(f"Error: Agencies file not found: {args.agencies_file}")
        return 1
    
    # Parse agencies
    agencies = []
    content = agencies_path.read_text()
    
    for line in content.split('\n'):
        if line.startswith('|') and 'https://' in line:
            parts = [p.strip() for p in line.split('|')]
            if len(parts) >= 4:
                name = parts[1]
                url = parts[2]
                folder = parts[3]
                if name and folder and name != 'Agency Name':
                    agencies.append({
                        'name': name,
                        'url': url,
                        'folder': folder
                    })
    
    print(f"Found {len(agencies)} agencies")
    
    # Filter if specific agency requested
    if args.agency:
        agencies = [a for a in agencies if a['folder'] == args.agency]
        if not agencies:
            print(f"Error: Agency '{args.agency}' not found")
            return 1
        print(f"Generating presentation for: {agencies[0]['name']}")
    
    # Initialize components
    print("\nInitializing presentation generator...")
    analyzer = AgencyContentAnalyzer()
    analyzer.load_network(args.graph_path)
    
    md_generator = MarkdownPresentationGenerator()
    ppt_converter = PowerPointConverter() if not args.md_only else None
    
    # Generate presentations
    print(f"\nGenerating presentations for {len(agencies)} agencies...\n")
    
    for i, agency in enumerate(agencies, 1):
        print(f"[{i}/{len(agencies)}] Processing: {agency['name']}")
        
        try:
            # Extract agency data
            mission_data = analyzer.extract_mission_overview(agency['folder'], agency['name'])
            mission_data['url'] = agency['url']
            
            services, keywords = analyzer.extract_services(agency['folder'])
            categories = analyzer.categorize_services(services, keywords)
            
            services_data = {
                'categories': categories,
                'keywords': keywords
            }
            
            # Generate markdown
            md_file = md_generator.generate_presentation(mission_data, services_data)
            
            # Convert to PowerPoint
            if ppt_converter:
                ppt_file = md_file.replace('.md', '.pptx')
                ppt_converter.convert_markdown_to_pptx(md_file, ppt_file)
            
            print(f"  ✓ Completed {agency['name']}\n")
            
        except Exception as e:
            print(f"  ✗ Error processing {agency['name']}: {e}\n")
            continue
    
    print("\n✅ Presentation generation complete!")
    print(f"Output directory: ../presentations/")
    
    return 0


if __name__ == "__main__":
    exit(main())
